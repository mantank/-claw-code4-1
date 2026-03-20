#!/usr/bin/env python3
"""
output-formatter.py - 把文本内容输出成不同格式
用法:
  python3 output-formatter.py excel "标题" "内容" > out.xlsx
  python3 output-formatter.py pptx "标题" "内容" > out.pptx
  python3 output-formatter.py pdf "标题" "内容" > out.pdf
  python3 output-formatter.py json "内容" > out.json
"""
import sys
import os

def to_excel(title, content):
    """输出Excel表格"""
    import pandas as pd
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment
    
    wb = Workbook()
    ws = wb.active
    ws.title = "数据"
    
    # 标题行
    ws['A1'] = title
    ws['A1'].font = Font(bold=True, size=14)
    ws.merge_cells('A1:Z1')
    
    # 解析内容（支持"指标|数值|备注"格式的表格）
    lines = content.strip().split('\n')
    for row_idx, line in enumerate(lines, start=3):
        if '|' in line:
            cells = [c.strip() for c in line.split('|')]
            for col_idx, val in enumerate(cells, start=1):
                cell = ws.cell(row=row_idx, column=col_idx, value=val)
                if row_idx == 3:
                    cell.font = Font(bold=True)
    
    path = f'/tmp/{title.replace(" ","_")}.xlsx'
    wb.save(path)
    return path

def to_pptx(title, content):
    """输出PPT"""
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.dml.color import RGBColor
    
    prs = Presentation()
    
    # 封面页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.title.text = title
    
    # 内容页（按段落分页）
    lines = content.strip().split('\n')
    current_slide = None
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if '|' in line:  # 表格行，跳到新页
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.title.text = title
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.text = '  |  '.join([c.strip() for c in line.split('|')])
        else:  # 普通文本
            if '---' in line:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.title.text = title
                continue
            if current_slide is None:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.title.text = title
                current_slide = slide
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
                tf = txBox.text_frame
                tf.word_wrap = True
                tf.text = line
            else:
                p = tf.add_paragraph()
                p.text = line
    
    path = f'/tmp/{title.replace(" ","_")}.pptx'
    prs.save(path)
    return path

def to_pdf(title, content):
    """输出PDF"""
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas
    
    path = f'/tmp/{title.replace(" ","_")}.pdf'
    c = canvas.Canvas(path, pagesize=A4)
    c.setFont("Helvetica-Bold", 20)
    c.drawString(50, 800, title)
    
    c.setFont("Helvetica", 11)
    y = 760
    for line in content.strip().split('\n'):
        if y < 50:
            c.showPage()
            y = 800
        if line.strip():
            c.drawString(50, y, line.strip()[:100])
            y -= 18
    
    c.save()
    return path

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)
    
    fmt = sys.argv[1]
    title = sys.argv[2] if len(sys.argv) > 2 else '未命名'
    content = sys.argv[3] if len(sys.argv) > 3 else ''
    
    # 如果内容是文件路径，读取文件
    if content.startswith('@'):
        with open(content[1:], 'r') as f:
            content = f.read()
    
    if fmt == 'excel':
        path = to_excel(title, content)
    elif fmt == 'pptx':
        path = to_pptx(title, content)
    elif fmt == 'pdf':
        path = to_pdf(title, content)
    else:
        print(f'不支持格式: {fmt}')
        sys.exit(1)
    
    print(f'✅ 已生成: {path}')
